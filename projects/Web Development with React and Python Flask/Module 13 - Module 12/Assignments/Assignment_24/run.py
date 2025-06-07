from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation
prs = Presentation()

# Title Slide
slide_layout = prs.slide_layouts[0]  # Title Slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "IELTS Speaking Test Platform"
subtitle.text = "Final Presentation\nTeam: Your Team Name\nDate: 2025-04-30"

# Introduction Slide
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Introduction"
content = slide.placeholders[1]
content.text = (
    "• Objective: Build an accessible IELTS Speaking Test platform.\n"
    "• Target Audience: IELTS candidates and administrators.\n"
)

# Platform Features Slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Platform Features"
content = slide.placeholders[1]
content.text = (
    "• User Authentication (JWT based)\n"
    "• Role-Based Dashboards (Admin, Test Taker)\n"
    "• Test-Taking Workflows with Timers\n"
    "• Real-Time Updates (WebSocket)\n"
    "• Accessibility Enhancements (ARIA, Keyboard Navigation)"
)

# Technical Highlights Slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Technical Highlights"
content = slide.placeholders[1]
content.text = (
    "• Frontend: React + Tailwind CSS\n"
    "• Backend: Flask + Flask-JWT-Extended\n"
    "• Database: PostgreSQL (SQLAlchemy ORM)\n"
    "• WebSocket: Real-time communication\n"
    "• Deployment: Docker & Azure"
)

# Demo Workflow Slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Demo Workflow"
content = slide.placeholders[1]
content.text = (
    "• Login as Test Taker\n"
    "• Start Speaking Test\n"
    "• Timer Starts and Questions Displayed\n"
    "• Submit Responses\n"
    "• Admin Reviews Submissions"
)

# Challenges and Solutions Slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Challenges and Solutions"
content = slide.placeholders[1]
content.text = (
    "• Challenge: Real-time updates with timer synchronization\n"
    "• Solution: Integrated WebSocket to broadcast timer events.\n"
    "• Challenge: Ensuring Accessibility Standards\n"
    "• Solution: Implemented ARIA roles, focus management, and tested with axe tools."
)

# Future Enhancements Slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Future Enhancements"
content = slide.placeholders[1]
content.text = (
    "• Add Speaking Evaluation using AI.\n"
    "• Multi-language Support.\n"
    "• Advanced Reporting for Admins.\n"
    "• Scalability Improvements (Kubernetes)."
)

# Conclusion Slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Conclusion"
content = slide.placeholders[1]
content.text = (
    "• Platform is ready for live demo.\n"
    "• Features tested and validated.\n"
    "• Thank you!"
)

# Save presentation
output_file = '/mnt/data/IELTS_Final_Presentation.pptx'
prs.save(output_file)

output_file
